"""
Generate docs/architecture-review.pptx for the PaperScout architecture review.

Re-runnable, idempotent. Reads no source files at runtime — all facts are
inlined (and were verified against the live source on 2026-05-17, branch main).

Usage:
    docs/.venv/bin/python docs/generate_architecture_review.py

Requires: python-pptx >= 1.0
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------- Constants ----------

OUT_PATH = Path(__file__).resolve().parent / "architecture-review.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PRIMARY = RGBColor(0x2C, 0x55, 0x82)   # blue — titles / boxes
ACCENT  = RGBColor(0x38, 0xA1, 0x69)   # green — pass / recommend
WARN    = RGBColor(0xDD, 0x6B, 0x20)   # orange — risk / store-only
DANGER  = RGBColor(0xC5, 0x3D, 0x3D)   # red — low quality / failures
BG      = RGBColor(0xF7, 0xFA, 0xFC)   # neutral background
TEXT    = RGBColor(0x1A, 0x20, 0x2C)   # body text
MUTED   = RGBColor(0x4A, 0x55, 0x68)   # muted captions
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RULE    = RGBColor(0xCB, 0xD5, 0xE0)   # divider lines

FOOTER  = "PaperScout · Architecture Review · 2026-05-17 · main @ 73966fc"

MONO_FONT = "Menlo"
SANS_FONT = "Helvetica Neue"


# ---------- Helpers ----------

def _set_text(tf, text, *, font=SANS_FONT, size=14, bold=False,
              color=TEXT, align=PP_ALIGN.LEFT):
    """Replace a text-frame's content with a single paragraph."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_title_bar(slide, title: str, eyebrow: str | None = None):
    """Top banner — blue accent rule + slide title + optional eyebrow tag."""
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.18)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY

    # Eyebrow (small uppercase tag, optional)
    if eyebrow:
        eb = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.32), Inches(12), Inches(0.3)
        )
        _set_text(eb.text_frame, eyebrow.upper(), size=11, bold=True,
                  color=PRIMARY)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.55) if eyebrow else Inches(0.45),
        Inches(12.1),
        Inches(0.7),
    )
    _set_text(title_box.text_frame, title, size=28, bold=True, color=TEXT)

    # Divider rule
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6),
        Inches(1.25) if eyebrow else Inches(1.15),
        Inches(12.1),
        Emu(12700),  # 1pt
    )
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE


def add_footer(slide):
    """Footer text bottom-left + page hint bottom-right (page filled later)."""
    fb = slide.shapes.add_textbox(
        Inches(0.6), Inches(7.1), Inches(11), Inches(0.3)
    )
    _set_text(fb.text_frame, FOOTER, size=9, color=MUTED)


def add_bullets(slide, items, *, left, top, width, height,
                font_size=14, line_spacing=1.2):
    """Bulleted list; items may be (label, body) tuples for bold label + body."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            label, body = item
            r1 = p.add_run()
            r1.text = "• " + label + ": "
            r1.font.name = SANS_FONT
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = PRIMARY
            r2 = p.add_run()
            r2.text = body
            r2.font.name = SANS_FONT
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = TEXT
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.name = SANS_FONT
            r.font.size = Pt(font_size)
            r.font.color.rgb = TEXT
    return tb


def add_code_block(slide, code, *, left, top, width, height, font_size=10):
    """Monospace text frame with subtle background for code/schema."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.color.rgb = RULE
    bg.line.width = Pt(0.5)
    bg.shadow.inherit = False

    tb = slide.shapes.add_textbox(
        left + Inches(0.12),
        top + Inches(0.08),
        width - Inches(0.24),
        height - Inches(0.16),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO_FONT
        r.font.size = Pt(font_size)
        r.font.color.rgb = TEXT


def add_box(slide, left, top, width, height, text, *,
            fill=PRIMARY, text_color=WHITE, font_size=12, bold=True,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, text, size=font_size, bold=bold, color=text_color,
              align=PP_ALIGN.CENTER)
    return shp


def add_arrow(slide, x1, y1, x2, y2, *, color=PRIMARY, weight=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # Arrow head
    try:
        from pptx.oxml.ns import qn
        ln = line.line._get_or_add_ln()
        tail = ln.find(qn("a:tailEnd"))
        if tail is None:
            from lxml import etree
            tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("h", "med")
    except Exception:
        pass
    return line


def add_table(slide, headers, rows, *, left, top, width, height,
              header_fill=PRIMARY, font_size=11):
    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, left, top, width, height
    ).table

    # Headers
    for c, h in enumerate(headers):
        cell = table_shape.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        _set_text(tf, h, size=font_size, bold=True, color=WHITE,
                  align=PP_ALIGN.LEFT)
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)

    # Rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table_shape.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else BG
            tf = cell.text_frame
            tf.word_wrap = True
            _set_text(tf, str(val), size=font_size, color=TEXT,
                      align=PP_ALIGN.LEFT)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return table_shape


# ---------- Slides ----------

def slide_01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background block
    bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Eyebrow
    eb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11), Inches(0.4))
    _set_text(eb.text_frame, "ARCHITECTURE REVIEW", size=14, bold=True,
              color=RGBColor(0xBE, 0xD9, 0xFF))

    # Main title
    t = s.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.3))
    _set_text(t.text_frame, "PaperScout", size=64, bold=True, color=WHITE)

    # Subtitle
    sub = s.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.6))
    _set_text(sub.text_frame,
              "Automated computer-vision paper collection, evaluation & ranking",
              size=22, color=RGBColor(0xE2, 0xE8, 0xF0))

    # Meta line
    meta = s.shapes.add_textbox(Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.4))
    _set_text(meta.text_frame,
              "Date: 2026-05-17    Branch: main    Commit: 73966fc    Stack: Next.js 16 · Prisma · PostgreSQL · Zod · Claude Code skills",
              size=12, color=RGBColor(0xCB, 0xD5, 0xE0))


def slide_02_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Problem & Goal", "Why this exists")

    add_bullets(s, [
        ("Problem", "≈100+ CV papers land daily across arXiv, OpenReview, "
                    "Hugging Face daily-papers. No single feed; manual triage "
                    "is impossible at this volume."),
        ("Goal", "Produce a daily, ranked, bilingual (en + zh-TW) digest of "
                  "~10 recommended papers with extracted key figures and "
                  "evidence-backed scores."),
        ("Approach", "Hybrid pipeline: Claude Code skills do collection + "
                     "LLM evaluation offline; a TypeScript CLI ingests + "
                     "dedups + ranks into Postgres; Next.js 16 serves a "
                     "read-only browse UI."),
    ], left=Inches(0.7), top=Inches(1.55), width=Inches(12),
       height=Inches(2.4), font_size=15)

    # Non-goals box
    add_box(s, Inches(0.7), Inches(4.6), Inches(12), Inches(2.1),
            "", fill=BG, shape=MSO_SHAPE.RECTANGLE)
    tb = s.shapes.add_textbox(Inches(0.95), Inches(4.7), Inches(11.5), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Non-goals (V1)"
    r.font.name = SANS_FONT; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = WARN
    for item in [
        "Real-time / streaming ingest — daily cadence is sufficient",
        "Full-text search across paper bodies — only abstracts + extracted snippets",
        "User accounts, comments, social features (PaperFeedback table exists but UI is anonymous)",
        "Multi-domain support — frozen to computer_vision in DailyRun.domain",
        "In-product paper authoring — read-only catalogue only",
    ]:
        p = tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run(); r.text = "  – " + item
        r.font.name = SANS_FONT; r.font.size = Pt(13)
        r.font.color.rgb = TEXT

    add_footer(s)


def slide_03_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "System at a Glance", "End-to-end data flow")

    # Sources cluster (top-left)
    src_x = Inches(0.7); src_y = Inches(1.9); src_w = Inches(1.7); src_h = Inches(0.55)
    add_box(s, src_x, src_y,            src_w, src_h, "arXiv",          fill=PRIMARY)
    add_box(s, src_x, src_y + Inches(0.7), src_w, src_h, "OpenReview",  fill=PRIMARY)
    add_box(s, src_x, src_y + Inches(1.4), src_w, src_h, "Hugging Face", fill=PRIMARY)

    # Collect skill
    col_x = Inches(3.0); col_y = Inches(2.6); col_w = Inches(2.0); col_h = Inches(0.85)
    collect = add_box(s, col_x, col_y, col_w, col_h,
                      "Collect skill\n(Claude Code)", fill=ACCENT, font_size=11)

    # candidates.json
    cand_x = Inches(5.4); cand_y = col_y; cand_w = Inches(1.6); cand_h = col_h
    add_box(s, cand_x, cand_y, cand_w, cand_h, "candidates.json",
            fill=MUTED, font_size=10, shape=MSO_SHAPE.PARALLELOGRAM)

    # Evaluate skill
    ev_x = Inches(7.4); ev_y = col_y; ev_w = Inches(2.0); ev_h = col_h
    add_box(s, ev_x, ev_y, ev_w, ev_h,
            "Evaluate skill\n(Stage 1 + 2)", fill=ACCENT, font_size=11)

    # evaluations.json
    ev_out_x = Inches(9.8); ev_out_y = col_y
    add_box(s, ev_out_x, ev_out_y, Inches(1.8), ev_h,
            "evaluations.json\n+ figures/*.png",
            fill=MUTED, font_size=10, shape=MSO_SHAPE.PARALLELOGRAM)

    # Phase 2.5 gate
    gate_y = Inches(4.2)
    add_box(s, Inches(7.4), gate_y, Inches(4.2), Inches(0.55),
            "Phase 2.5 harness gate  (npm run prompt:check)",
            fill=WARN, font_size=11)

    # Ingest
    ing_y = Inches(5.2)
    add_box(s, Inches(5.4), ing_y, Inches(3.0), Inches(0.7),
            "Ingest CLI  scripts/ingest.ts",
            fill=PRIMARY, font_size=12)

    # DB
    add_box(s, Inches(8.7), ing_y, Inches(1.6), Inches(0.7),
            "PostgreSQL\n(Prisma)", fill=DANGER, font_size=11,
            shape=MSO_SHAPE.CAN)

    # Next.js UI
    add_box(s, Inches(10.5), ing_y, Inches(2.1), Inches(0.7),
            "Next.js 16 UI\n/, /papers/[id], /library, /runs/[id]",
            fill=PRIMARY, font_size=10)

    # Arrows: sources → collect
    for y_off in (0.0, 0.7, 1.4):
        add_arrow(s,
                  src_x + src_w, src_y + Inches(0.27 + y_off),
                  col_x,         col_y + Inches(0.43))
    # collect → candidates
    add_arrow(s, col_x + col_w, col_y + Inches(0.43),
              cand_x, cand_y + Inches(0.43))
    # candidates → evaluate
    add_arrow(s, cand_x + cand_w, cand_y + Inches(0.43),
              ev_x, ev_y + Inches(0.43))
    # evaluate → evaluations
    add_arrow(s, ev_x + ev_w, ev_y + Inches(0.43),
              ev_out_x, ev_out_y + Inches(0.43))
    # evaluations → gate
    add_arrow(s, ev_out_x + Inches(0.9), ev_out_y + ev_h,
              Inches(9.5), gate_y)
    # gate → ingest
    add_arrow(s, Inches(9.5), gate_y + Inches(0.55),
              Inches(6.9), ing_y)
    # ingest → DB → UI
    add_arrow(s, Inches(8.4), ing_y + Inches(0.35),
              Inches(8.7), ing_y + Inches(0.35))
    add_arrow(s, Inches(10.3), ing_y + Inches(0.35),
              Inches(10.5), ing_y + Inches(0.35))

    # Divider: offline vs online
    div = s.shapes.add_connector(1, Inches(0.7), Inches(4.95),
                                 Inches(12.7), Inches(4.95))
    div.line.color.rgb = RULE
    div.line.width = Pt(0.75)
    div.line.dash_style = 7  # dash

    lbl1 = s.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(5), Inches(0.3))
    _set_text(lbl1.text_frame, "OFFLINE — Claude Code skills + harness",
              size=10, bold=True, color=MUTED)
    lbl2 = s.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(5), Inches(0.3))
    _set_text(lbl2.text_frame, "ONLINE — TypeScript CLI + Next.js runtime",
              size=10, bold=True, color=MUTED)

    add_footer(s)


def slide_04_stack(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Tech Stack", "Pinned versions, package.json")

    headers = ["Layer", "Component", "Version", "Role"]
    rows = [
        ["Frontend",  "next",              "16.2.5",  "App-Router SSR, edge-ready"],
        ["Frontend",  "react / react-dom", "19.2.4",  "RSC + client components"],
        ["Frontend",  "tailwindcss",       "^4",      "Utility CSS via @tailwindcss/postcss"],
        ["Frontend",  "@radix-ui/*",       "~1.x",    "Dialog, Select, Slider, Progress"],
        ["Backend",   "@prisma/client",    "^6.19.3", "Postgres ORM"],
        ["Backend",   "zod",               "^4.4.3",  "Schema validation (candidate / evaluation)"],
        ["Backend",   "fast-xml-parser",   "^5.7.3",  "arXiv Atom feed parsing"],
        ["Backend",   "pino",              "^10.3.1", "Structured logging"],
        ["Tooling",   "tsx",               "^4.21.0", "Runs scripts/ + Prisma CLI in TS"],
        ["Testing",   "vitest",            "^4.1.5",  "Unit + integration (real Postgres)"],
        ["Testing",   "@playwright/test",  "^1.59.1", "Browser E2E"],
        ["Dedup",     "fastest-levenshtein", "^1.0.16", "Fuzzy title matching"],
    ]
    add_table(s, headers, rows,
              left=Inches(0.7), top=Inches(1.55),
              width=Inches(12), height=Inches(5.0),
              font_size=11)

    add_footer(s)


def slide_05_repo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Repo Layout", "Top-level + src/")

    tree = """paper-collection-system/
├── src/
│   ├── app/                       # Next.js 16 App Router
│   │   ├── page.tsx               # Home feed (paginated, locale-aware)
│   │   ├── library/, papers/[id]/, runs/[id]/
│   │   └── api/{locale,runs,papers/[id]/figure}/
│   ├── server/                    # 'server-only' boundary
│   │   ├── schema/                # candidate.ts, evaluation.ts (Zod)
│   │   ├── sources/               # arxiv.ts, openreview.ts, huggingface.ts, index.ts
│   │   ├── dedup/                 # matcher.ts, fingerprint.ts, normalize.ts
│   │   ├── repos/                 # papers, evaluations, runs, runResults, trends, …
│   │   ├── pipeline/              # collect.ts, persist.ts, runner.ts (reserved)
│   │   └── lib/                   # select-evaluation.ts (best-eval picker)
│   ├── components/ui/             # Radix + Tailwind primitives
│   ├── lib/                       # client utils: locale, format, db
│   └── i18n/                      # en.ts ↔ zh-TW.ts (mirror shape)
├── prisma/                        # schema.prisma + migrations + seed.ts
├── scripts/
│   ├── ingest.ts                  # Phase 3 ingest CLI
│   ├── validate-{candidates,evaluations}.ts
│   └── prompt-eval/               # Phase 2.5 harness + F1..F5 fixtures
├── tests/
│   ├── unit/{sources,dedup,prompt-eval}/
│   └── integration/{ingest,collect-persist}.test.ts   # real Postgres
├── .claude/skills/
│   ├── collect-papers/SKILL.md    # Phase 1 spec
│   └── evaluate-papers/SKILL.md   # Phase 2 spec
└── data/
    ├── runs/<YYYY-MM-DD-HHMM>/    # candidates.json + evaluations.json + figures/
    └── sample/                    # reference shape for the two artefacts"""
    add_code_block(s, tree,
                   left=Inches(0.7), top=Inches(1.55),
                   width=Inches(12), height=Inches(5.3), font_size=10)
    add_footer(s)


def slide_06_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Pipeline Phases", "Three phases + a gating harness")

    # Swim-lane: 5 vertical phase boxes, top row = name, middle = artefact, bottom = exec mode
    lanes = [
        ("Phase 1\nCollect",       "candidates.json", "Claude Code skill\n(offline)",      ACCENT),
        ("Phase 2a\nAbstract",     "scores + tags",   "Claude Code skill\n(LLM eval)",     ACCENT),
        ("Phase 2b\nFull-PDF top 15", "PDF + figures", "Claude Code skill\n(LLM + Bash)",  ACCENT),
        ("Phase 2.5\nHarness gate","pass / fail",     "tsx CLI\n(local check)",            WARN),
        ("Phase 3\nIngest",        "Postgres rows",   "tsx CLI\n(idempotent)",             PRIMARY),
    ]
    lane_w = Inches(2.35)
    lane_h_top = Inches(0.8)
    lane_h_mid = Inches(0.7)
    lane_h_bot = Inches(0.7)
    gap = Inches(0.18)
    start_x = Inches(0.7)
    base_y = Inches(1.7)

    for i, (name, artefact, mode, color) in enumerate(lanes):
        x = start_x + i * (lane_w + gap)
        add_box(s, x, base_y, lane_w, lane_h_top, name, fill=color, font_size=13)
        add_box(s, x, base_y + lane_h_top + Inches(0.05),
                lane_w, lane_h_mid, artefact, fill=BG,
                text_color=TEXT, font_size=11, bold=False,
                shape=MSO_SHAPE.RECTANGLE)
        add_box(s, x, base_y + lane_h_top + lane_h_mid + Inches(0.1),
                lane_w, lane_h_bot, mode, fill=WHITE,
                text_color=MUTED, font_size=10, bold=False,
                shape=MSO_SHAPE.RECTANGLE)

        if i < len(lanes) - 1:
            ax1 = x + lane_w
            ay  = base_y + Inches(0.4)
            ax2 = ax1 + gap
            add_arrow(s, ax1, ay, ax2, ay)

    # Bottom annotations
    notes_y = Inches(4.6)
    add_bullets(s, [
        ("Handoffs",
         "Phases 1↔2 share a run-dir; 2↔2.5 share evaluations.json; "
         "2.5↔3 only proceeds if gate passes (≥4/5 fixtures within bounds)."),
        ("Idempotency",
         "Phase 3 dedups by SHA-256 of the SKILL.md (llmPromptVersion) "
         "and refuses re-ingest of the same ingestSourceDir without --force."),
        ("Fault isolation",
         "Phase 1 uses Promise.allSettled — one source down (e.g. arXiv 503) "
         "does not block the run; failures captured in meta.errors[]."),
        ("Cadence",
         "Designed for ~daily runs (~30 candidates → ~15 PDF reads → ~10 RECOMMEND). "
         "No streaming; no per-paper async pipeline."),
    ], left=Inches(0.7), top=notes_y, width=Inches(12), height=Inches(2.4),
       font_size=12)

    add_footer(s)


def slide_07_collection_sources(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Phase 1 — Collection Sources",
                  "src/server/sources/{arxiv,openreview,huggingface}.ts")

    headers = ["Source", "Quota", "Endpoint", "Parser",
               "Per-source notes"]
    rows = [
        [
            "arXiv",
            "15 / 30",
            "export.arxiv.org/api/query\n?search_query=cat:cs.CV\n&sortBy=submittedDate",
            "Atom XML\nfast-xml-parser",
            "3 s self-rate-limit; extracts arXiv ID from URL; "
            "scans abstract for GitHub URLs → codeUrls[]",
        ],
        [
            "OpenReview",
            "5 / 30",
            "api2.openreview.net/notes/search\n?term=computer+vision\n&sort=cdate:desc",
            "JSON\n(nested content.value)",
            "Created / published dates come as epoch-ms → "
            "normalised to ISO-8601 YYYY-MM-DD",
        ],
        [
            "Hugging Face",
            "10 / 30",
            "huggingface.co/api/daily_papers",
            "JSON\n(post-filtered)",
            "CV keyword allow-list (vision / image / video / "
            "diffusion / …); detects arXiv slug for cross-source dedup hint",
        ],
    ]
    add_table(s, headers, rows,
              left=Inches(0.55), top=Inches(1.55),
              width=Inches(12.3), height=Inches(4.4),
              font_size=11)

    # Bottom call-out
    add_box(s, Inches(0.55), Inches(6.15), Inches(12.3), Inches(0.75),
            "Common HTTP wrapper: src/server/sources/http.ts — AbortController "
            "timeout, dep-injectable fetch (testable). Aggregator: "
            "src/server/sources/index.ts.",
            fill=BG, text_color=TEXT, font_size=11, bold=False,
            shape=MSO_SHAPE.RECTANGLE)

    add_footer(s)


def slide_08_aggregation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Phase 1 — Aggregation & In-Batch Dedup",
                  "src/server/sources/index.ts")

    add_bullets(s, [
        ("Parallel fetch",
         "Promise.allSettled([arxiv, openreview, huggingface]) — "
         "per-source failures are logged to meta.errors[] but never abort the run."),
        ("Within-batch dedup keys",
         "Exact arXiv ID, exact OpenReview ID, then normalized title "
         "(see src/server/dedup/normalize.ts)."),
        ("Source priority on collision",
         "ARXIV (0) > OPENREVIEW (1) > HUGGINGFACE (2). Winner becomes "
         "the primary record; loser is folded into additionalSources[]."),
        ("Two-pass quota enforcement",
         "(1) Per-source quotas fill greedily; (2) any deficit is back-"
         "filled from leftovers across all sources until the ~30 target is hit."),
        ("Output contract",
         "data/runs/<YYYY-MM-DD-HHMM>/candidates.json conforming to "
         "CandidatesFileSchema (Zod). Validated post-write via "
         "`npm run validate:candidates <path>`."),
    ], left=Inches(0.7), top=Inches(1.55), width=Inches(12),
       height=Inches(3.6), font_size=14)

    # Mini diagram showing the two-pass quota fill
    diag_y = Inches(5.4)
    add_box(s, Inches(0.7),  diag_y, Inches(2.4), Inches(0.6),
            "Pass 1: per-source quotas",
            fill=PRIMARY, font_size=11)
    add_arrow(s, Inches(3.1), diag_y + Inches(0.3),
              Inches(3.5), diag_y + Inches(0.3))
    add_box(s, Inches(3.5), diag_y, Inches(2.4), Inches(0.6),
            "Within-batch dedup",
            fill=PRIMARY, font_size=11)
    add_arrow(s, Inches(5.9), diag_y + Inches(0.3),
              Inches(6.3), diag_y + Inches(0.3))
    add_box(s, Inches(6.3), diag_y, Inches(2.4), Inches(0.6),
            "Pass 2: back-fill deficit",
            fill=PRIMARY, font_size=11)
    add_arrow(s, Inches(8.7), diag_y + Inches(0.3),
              Inches(9.1), diag_y + Inches(0.3))
    add_box(s, Inches(9.1), diag_y, Inches(3.4), Inches(0.6),
            "Write candidates.json",
            fill=ACCENT, font_size=11)

    add_footer(s)


def slide_09_candidate_schema(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "CandidateRecord Schema",
                  "src/server/schema/candidate.ts (Zod)")

    code = """export const SourceEnum = z.enum(['ARXIV', 'OPENREVIEW', 'HUGGINGFACE']);

export const AdditionalSourceSchema = z.object({
  source: SourceEnum,
  sourceUrl: z.string().url(),
  sourcePaperId: z.string().min(1),
});

export const CandidateSchema = z.object({
  title:             z.string().min(1),
  authors:           z.array(z.string().min(1)).min(1),
  abstract:          z.string().nullable(),
  venue:             z.string().nullable(),
  publishedDate:     z.string().regex(/^\\d{4}-\\d{2}-\\d{2}$/, '…ISO-8601…'),
  sourceUrl:         z.string().url(),
  pdfUrl:            z.string().url().nullable(),
  sourcePaperId:     z.string().min(1).nullable(),
  source:            SourceEnum,
  codeUrls:          z.array(z.string().url()).default([]),
  additionalSources: z.array(AdditionalSourceSchema).default([]),
});

export const CandidatesFileSchema = z.array(CandidateSchema);
export type Candidate = z.infer<typeof CandidateSchema>;"""
    add_code_block(s, code, left=Inches(0.55), top=Inches(1.55),
                   width=Inches(8.2), height=Inches(5.3), font_size=10)

    # Right-side notes
    add_bullets(s, [
        ("Why a separate schema file?",
         "Pure Zod — no DB / env imports — so it can be loaded by either the "
         "ingest CLI or the validate-candidates script with zero ceremony."),
        ("Cross-source bridge",
         "additionalSources[] carries hints from in-batch dedup so the "
         "Phase 3 cross-run matcher can recognise the same paper across runs."),
        ("Strictness",
         "Every string field is min-1; publishedDate is regex-validated; "
         "URLs are .url() validated."),
        ("File location",
         "src/server/schema/candidate.ts (33 LoC)"),
    ], left=Inches(8.9), top=Inches(1.55),
       width=Inches(4.0), height=Inches(5.3), font_size=11)

    add_footer(s)


def slide_10_abstract_screening(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Phase 2a — Abstract Screening",
                  ".claude/skills/evaluate-papers/SKILL.md")

    # 5 score-dimension chips
    chips = [
        ("Novelty",                   "0 – 25", PRIMARY),
        ("Methodological\nRigor",     "0 – 25", PRIMARY),
        ("Experimental\nQuality",     "0 – 20", PRIMARY),
        ("Venue\nCredibility",        "0 – 15", PRIMARY),
        ("Author / Inst\nReputation", "0 – 15", PRIMARY),
    ]
    chip_w = Inches(2.35); chip_h = Inches(1.0)
    chip_y = Inches(1.55)
    for i, (name, rng, color) in enumerate(chips):
        x = Inches(0.55) + i * (chip_w + Inches(0.15))
        add_box(s, x, chip_y, chip_w, chip_h,
                f"{name}\n{rng}", fill=color, font_size=12)

    # Total + thresholds row
    add_box(s, Inches(0.55), chip_y + Inches(1.2), Inches(3.5), Inches(0.7),
            "TOTAL  =  Σ five dimensions  ∈ 0..100",
            fill=BG, text_color=TEXT, font_size=13,
            shape=MSO_SHAPE.RECTANGLE)
    add_box(s, Inches(4.2), chip_y + Inches(1.2), Inches(2.8), Inches(0.7),
            "≥ 65  →  RECOMMEND", fill=ACCENT, font_size=13)
    add_box(s, Inches(7.15), chip_y + Inches(1.2), Inches(2.8), Inches(0.7),
            "50 – 64  →  STORE_ONLY", fill=WARN, font_size=13)
    add_box(s, Inches(10.1), chip_y + Inches(1.2), Inches(2.8), Inches(0.7),
            "< 50  →  LOW_QUALITY", fill=DANGER, font_size=13)

    # Detail bullets
    add_bullets(s, [
        ("Inputs",
         "Title, authors, abstract, venue, source metadata only — no PDF fetch yet."),
        ("Coverage",
         "All ~30 candidates are scored. Schema enforces "
         "scores.total = Σ five dims via z.refine (otherwise the file fails validation)."),
        ("Bilingual outputs",
         "summary, recommendationReason, rankingExplanation are all "
         "{ en, 'zh-TW' } — both locales required, min-1 char each."),
        ("Tags",
         "2 – 5 lowercase English keywords (e.g. 'vision-transformer', "
         "'segmentation'). Tags are NOT bilingual — they are stable indexes."),
        ("Stage marker",
         "evaluationStage = 'ABSTRACT_SCREENING'; pdfAnalysisStatus, "
         "keyContribution, methodologySummary, strengths, weaknesses, figure "
         "MUST be null at this stage (enforced by superRefine)."),
    ], left=Inches(0.55), top=Inches(4.5), width=Inches(12.3),
       height=Inches(2.5), font_size=12)

    add_footer(s)


def slide_11_full_pdf(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Phase 2b — Full-PDF Read (top 15)",
                  "PDF truncate + refine + figure pull")

    # Horizontal pipeline diagram
    py = Inches(1.7)
    steps = [
        ("Rank Stage-1\nby total → top 15",  PRIMARY),
        ("curl\n32 MB / 60 s cap",           PRIMARY),
        ("pdftotext +\nappendix marker",     PRIMARY),
        ("qpdf --pages 1..N\ntruncate",      PRIMARY),
        ("Read tool\n(LLM ingest)",          ACCENT),
        ("Refine scores +\nfill narratives", ACCENT),
    ]
    step_w = Inches(1.95); step_h = Inches(0.95)
    for i, (label, color) in enumerate(steps):
        x = Inches(0.55) + i * (step_w + Inches(0.1))
        add_box(s, x, py, step_w, step_h, label, fill=color, font_size=11)
        if i < len(steps) - 1:
            add_arrow(s, x + step_w, py + Inches(0.47),
                      x + step_w + Inches(0.1), py + Inches(0.47))

    # Bullet details
    add_bullets(s, [
        ("Prereqs (Bash)",
         "poppler (pdftotext, pdftocairo) and qpdf must be on PATH. "
         "The skill aborts that PDF if either is missing."),
        ("Truncation heuristic",
         "Scan text for 'Appendix' or 'Supplementary Material'; "
         "qpdf --pages 1-<n> slices off everything from that page on. "
         "Fallback: keep entire PDF (n/a or marker on page 1)."),
        ("Refined fields (bilingual, all required when SUCCESS)",
         "keyContribution (1–2 sent.), methodologySummary (2–3 sent.), "
         "strengths[] (3–5 per locale), weaknesses[] (2–4 per locale). "
         "List lengths AND order must mirror across en ↔ zh-TW."),
        ("Stage marker on failure",
         "pdfAnalysisStatus = 'FAILED' (parse error) | 'UNAVAILABLE' "
         "(fetch / size cap). Stage-1 record is preserved verbatim — no "
         "refinement claimed without evidence."),
    ], left=Inches(0.55), top=Inches(3.05), width=Inches(12.3),
       height=Inches(3.9), font_size=12)

    add_footer(s)


def slide_12_figure_extraction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Phase 2b — Figure Extraction",
                  "Two-pass render via pdftocairo")

    # Two-pass diagram
    py = Inches(1.7)
    add_box(s, Inches(0.55), py, Inches(2.5), Inches(0.95),
            "Pass 1: preview\npdftocairo -png -r 72",
            fill=PRIMARY, font_size=11)
    add_arrow(s, Inches(3.05), py + Inches(0.47),
              Inches(3.25), py + Inches(0.47))
    add_box(s, Inches(3.25), py, Inches(2.8), Inches(0.95),
            "LLM picks figure +\nestimates bbox (x,y,w,h)\nas frac of page",
            fill=ACCENT, font_size=11)
    add_arrow(s, Inches(6.05), py + Inches(0.47),
              Inches(6.25), py + Inches(0.47))
    add_box(s, Inches(6.25), py, Inches(3.1), Inches(0.95),
            "Pass 2: crop\npdftocairo -png -r 150\n-x -y -W -H",
            fill=PRIMARY, font_size=11)
    add_arrow(s, Inches(9.35), py + Inches(0.47),
              Inches(9.55), py + Inches(0.47))
    add_box(s, Inches(9.55), py, Inches(3.3), Inches(0.95),
            "figures/<safe-id>.png\n+ FigureSchema entry",
            fill=ACCENT, font_size=11)

    add_bullets(s, [
        ("Preference order",
         "Architecture diagram  >  main result figure  >  Figure 1 / teaser."),
        ("Caption rules",
         "caption.en = verbatim from PDF (trim 'Figure N:' prefix), ≤ 240 chars. "
         "caption.zh-TW = faithful Traditional-Chinese translation, ≤ 240 chars. "
         "Both enforced by Zod max-length."),
        ("Storage",
         "renderedPath is RELATIVE to the run dir (e.g. 'figures/2605.12345.png'). "
         "Ingest reads those bytes and stores them in PaperFigure.imageBytes (BYTEA)."),
        ("Invariant",
         "figure may only be non-null when pdfAnalysisStatus = 'SUCCESS' — "
         "enforced by EvaluationSchema.superRefine."),
    ], left=Inches(0.55), top=Inches(3.0), width=Inches(12.3),
       height=Inches(3.9), font_size=12)

    add_footer(s)


def slide_13_eval_schema(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "EvaluationRecord Schema",
                  "src/server/schema/evaluation.ts — superRefine invariants")

    code = """const localizedString = (opts = {}) => z.object({
  en:        opts.max ? z.string().min(1).max(opts.max) : z.string().min(1),
  'zh-TW':   opts.max ? z.string().min(1).max(opts.max) : z.string().min(1),
}).strict();

export const ScoresSchema = z.object({
  novelty:                      z.number().int().min(0).max(25),
  methodologicalRigor:          z.number().int().min(0).max(25),
  experimentalQuality:          z.number().int().min(0).max(20),
  venueSourceCredibility:       z.number().int().min(0).max(15),
  authorInstitutionReputation:  z.number().int().min(0).max(15),
  total:                        z.number().int().min(0).max(100),
}).refine(s => s.total === <sum of 5>, { message: '…' });

export const EvaluationSchema = z.object({
  joinKey:                { source, sourcePaperId },
  evaluationStage:        z.enum(['ABSTRACT_SCREENING', 'FULL_PDF']),
  scores:                 ScoresSchema,
  summary:                localizedString(),
  recommendationReason:   localizedString(),
  keyContribution:        localizedString().nullable(),
  methodologySummary:     localizedString().nullable(),
  strengths:              localizedStringList.nullable(),
  weaknesses:             localizedStringList.nullable(),
  tags:                   z.array(z.string().min(1)).default([]),
  rankingExplanation:     localizedString(),
  recommendationDecision: z.enum(['RECOMMEND','STORE_ONLY','LOW_QUALITY']),
  pdfAnalysisStatus:      z.enum(['SUCCESS','FAILED','UNAVAILABLE']).nullable(),
  figure:                 FigureSchema.nullable().default(null),
}).superRefine(/* cross-field invariants — see next box */)"""
    add_code_block(s, code, left=Inches(0.55), top=Inches(1.55),
                   width=Inches(8.4), height=Inches(5.2), font_size=9)

    add_bullets(s, [
        ("Cross-field invariants (lines 84–139)",
         ""),
        ("stage == FULL_PDF",
         "pdfAnalysisStatus required."),
        ("pdfAnalysisStatus == SUCCESS",
         "keyContribution, methodologySummary, strengths, weaknesses "
         "all required (≥ 1 per locale)."),
        ("stage == ABSTRACT_SCREENING",
         "pdfAnalysisStatus MUST be null; figure MUST be null."),
        ("figure",
         "may only exist when pdfAnalysisStatus == SUCCESS — irrespective "
         "of stage."),
    ], left=Inches(9.1), top=Inches(1.55), width=Inches(3.8),
       height=Inches(5.2), font_size=11)

    add_footer(s)


def slide_14_bilingual(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Bilingual-First Design",
                  "Schema-enforced en + zh-TW across the stack")

    # Top: schema declaration
    code = """// src/server/schema/evaluation.ts
export const SUPPORTED_LOCALES = ['en', 'zh-TW'] as const;
type LocalizedString     = { en: string; 'zh-TW': string };
type LocalizedStringList = { en: string[]; 'zh-TW': string[] };

// .strict() => extra locales are a validation error
// list lengths and order MUST mirror across en ↔ zh-TW"""
    add_code_block(s, code, left=Inches(0.55), top=Inches(1.55),
                   width=Inches(12.3), height=Inches(1.55), font_size=11)

    # Two columns: backend rules / frontend rendering
    add_box(s, Inches(0.55), Inches(3.3), Inches(6.0), Inches(0.5),
            "Backend rules", fill=PRIMARY, font_size=13)
    add_bullets(s, [
        "Every narrative field of EvaluationRecord is bilingual",
        "strengths[].en[i] ↔ strengths['zh-TW'][i] — index-aligned",
        "1:1 sentence alignment within each pair",
        "Technical terms (ViT, ImageNet-1k) kept verbatim in zh-TW",
        "Tags are English-only (stable index keys, not display strings)",
        "DB columns store the raw JSON (`Json?` in PaperEvaluation)",
    ], left=Inches(0.6), top=Inches(3.85), width=Inches(6),
       height=Inches(2.9), font_size=11)

    add_box(s, Inches(6.85), Inches(3.3), Inches(6.0), Inches(0.5),
            "Frontend rendering", fill=PRIMARY, font_size=13)
    add_bullets(s, [
        "Cookie `locale` ∈ {en, zh-TW}; API: /api/locale (route handler)",
        "src/lib/locale.ts exports pickLocalized(value, locale)",
        "Messages catalogs: src/i18n/en.ts ↔ src/i18n/zh-TW.ts",
        "Catalogs must mirror shape (keys, depth) — values differ",
        "common.{sources,decisions} marked `as const` for typed keys",
        "Body fields (summary, etc.) selected at render time, not at API",
    ], left=Inches(6.9), top=Inches(3.85), width=Inches(6),
       height=Inches(2.9), font_size=11)

    add_footer(s)


def slide_15_prompt_harness(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Phase 2.5 — Prompt Harness",
                  "scripts/prompt-eval/ — gates evaluation quality before ingest")

    add_bullets(s, [
        ("Why",
         "LLM outputs drift between model versions and prompt edits. "
         "Without a gate, regressions silently propagate to Postgres."),
        ("Fixtures (F1 … F5)",
         "5 hand-curated CV papers spanning the quality spectrum, each with "
         "metadata.json (CandidateSchema + _fixture block) and bounds.json "
         "(per-dimension min/max + allowed enums)."),
        ("Gate logic (check-evaluations.ts, ~lines 124–142)",
         "PASS iff:  all records pass EvaluationSchema  AND  ≥ 4 of 5 fixtures "
         "within bounds  AND  every joinKey resolves  AND  fixture coverage "
         "is complete (no missing / duplicate / unexpected)."),
        ("Coarse ranking flags (informational)",
         "f1InTop2 (highest-quality fixture in top 2), f5InBottom2 "
         "(lowest-quality fixture in bottom 2), f4NotLast (mid-tier not last)."),
        ("Versioning",
         "Phase 3 records SHA-256(SKILL.md)[0:12] as llmPromptVersion on "
         "every PaperEvaluation — provenance survives across runs."),
    ], left=Inches(0.55), top=Inches(1.55), width=Inches(12.3),
       height=Inches(3.6), font_size=12)

    # Scripts row
    add_box(s, Inches(0.55), Inches(5.5), Inches(4.0), Inches(0.7),
            "npm run prompt:fixtures\n→ build reference run from F1..F5",
            fill=PRIMARY, font_size=10)
    add_box(s, Inches(4.75), Inches(5.5), Inches(4.0), Inches(0.7),
            "npm run prompt:check  <run-dir>\n→ exits 1 on gate failure",
            fill=WARN, font_size=10)
    add_box(s, Inches(8.95), Inches(5.5), Inches(3.9), Inches(0.7),
            "npm run prompt:normalize\n→ post-ingest stat normalization",
            fill=PRIMARY, font_size=10)

    add_footer(s)


def slide_16_ingest(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Phase 3 — Ingest CLI",
                  "scripts/ingest.ts  (npm run ingest <run-dir>)")

    add_bullets(s, [
        ("1. Validate",
         "Zod-parse candidates.json + evaluations.json; every evaluation "
         "joinKey must resolve to a candidate (else fail-fast)."),
        ("2. Idempotency",
         "DailyRun is keyed by ingestSourceDir (UNIQUE). If already ingested, "
         "exit cleanly; `--force` cascade-deletes the prior DailyRun first."),
        ("3. Prompt version",
         "SHA-256(.claude/skills/evaluate-papers/SKILL.md)[0:12] — "
         "stamped on every new PaperEvaluation as llmPromptVersion."),
        ("4. Cross-run dedup",
         "src/server/dedup/matcher.ts findMatch() — see Slide 17."),
        ("5. Persist candidates",
         "Upsert Paper (by duplicateFingerprint), insert PaperSource for "
         "primary + each additionalSource."),
        ("6. Persist evaluations + figures",
         "Insert PaperEvaluation (unique on (paperId, runId, stage)); "
         "read PNG bytes from disk, insert PaperFigure."),
        ("7. Rank & flag",
         "src/server/lib/select-evaluation.ts picks the 'best' eval per paper "
         "(SUCCESS > ABSTRACT > FAILED/UNAVAILABLE), then orders by totalScore. "
         "Writes finalRank + isRecommended into PaperRunResult."),
        ("8. Finalize",
         "DailyRun.status = COMPLETED, completedAt = now(). "
         "Prints summary (papers / evals / figures / wall-time)."),
    ], left=Inches(0.55), top=Inches(1.55), width=Inches(12.3),
       height=Inches(5.4), font_size=12)

    add_footer(s)


def slide_17_dedup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Cross-Run Dedup Strategy",
                  "src/server/dedup/{matcher,fingerprint,normalize}.ts → PaperDuplicate")

    headers = ["#", "MatchMethod", "Confidence", "Signal",
               "When it triggers"]
    rows = [
        ["1", "ARXIV_ID",         "1.00",
         "PaperSource where source=ARXIV AND sourcePaperId == new id",
         "Same arXiv preprint reposted or seen on another source"],
        ["2", "OPENREVIEW_ID",    "1.00",
         "PaperSource where source=OPENREVIEW AND sourcePaperId == new id",
         "Same OpenReview submission re-encountered"],
        ["3", "PDF_URL",          "0.95",
         "Exact-match PaperSource.pdfUrl (canonicalised)",
         "Mirror / alternate landing page with same PDF"],
        ["4", "SOURCE_URL",       "0.90",
         "Exact-match PaperSource.sourceUrl (canonicalised)",
         "Same listing URL pulled by a different fetcher"],
        ["5", "NORMALIZED_TITLE", "0.85",
         "Paper.normalizedTitle == normalize(new title)",
         "Title-collision after whitespace/case/punct normalisation"],
        ["6", "FUZZY_TITLE",      "var.",
         "fastest-levenshtein over normalised titles ≥ threshold",
         "Camera-ready vs preprint title drift, minor typo edits"],
    ]
    add_table(s, headers, rows,
              left=Inches(0.4), top=Inches(1.55),
              width=Inches(12.5), height=Inches(4.4), font_size=10)

    add_bullets(s, [
        ("Order matters",
         "matcher.ts tries methods 1→6 in order; first hit wins. Each match "
         "creates a PaperDuplicate row (canonicalPaperId, duplicatePaperId, "
         "matchMethod, confidence) — full audit trail."),
        ("Deterministic fingerprint",
         "Paper.duplicateFingerprint is the canonical write-once dedup key "
         "(unique index), derived from (primarySource, sourcePaperId | "
         "normalisedTitle)."),
    ], left=Inches(0.55), top=Inches(6.05), width=Inches(12.3),
       height=Inches(1.0), font_size=11)

    add_footer(s)


def slide_18_db_schema(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Database Schema",
                  "prisma/schema.prisma — 10 models, PostgreSQL")

    # Center: Paper
    paper = add_box(s, Inches(5.5), Inches(3.45), Inches(2.3), Inches(0.85),
                    "Paper\n(id, title, normalizedTitle,\nduplicateFingerprint)",
                    fill=PRIMARY, font_size=10)

    # Satellites — placed around Paper
    nodes = [
        # (label, x, y, w, h, color, edge-label)
        ("PaperSource\n(primary + extras)",     Inches(1.9),  Inches(1.85), Inches(2.2), Inches(0.7), ACCENT, "1..N"),
        ("PaperFigure\n(BYTEA imageBytes)",     Inches(1.9),  Inches(3.45), Inches(2.2), Inches(0.7), ACCENT, "0..1"),
        ("PaperTag\n(LLM / USER)",              Inches(1.9),  Inches(5.05), Inches(2.2), Inches(0.7), ACCENT, "0..N"),
        ("PaperCodeLink\n(github / pwc / …)",   Inches(9.3),  Inches(1.85), Inches(2.5), Inches(0.7), ACCENT, "0..N"),
        ("PaperEvaluation\n(5 score cols + JSON narratives)", Inches(9.3),  Inches(3.45), Inches(2.9), Inches(0.85), PRIMARY, "1..N"),
        ("PaperDuplicate\n(canonical ↔ dup, matchMethod)", Inches(9.3),  Inches(5.05), Inches(2.9), Inches(0.85), WARN, "self-ref"),
        ("PaperFeedback\n(star, comment)",      Inches(1.9),  Inches(6.25), Inches(2.2), Inches(0.7), MUTED, "0..N"),
    ]
    for label, x, y, w, h, color, _ in nodes:
        add_box(s, x, y, w, h, label, fill=color, font_size=10)
        # straight-ish connector to Paper center
        px = x + w if x < Inches(5) else x
        py = y + h / 2
        cx = Inches(5.5) if x < Inches(5) else Inches(5.5) + Inches(2.3)
        cy = Inches(3.45) + Inches(0.425)
        add_arrow(s, px, py, cx, cy, color=RULE, weight=1.0)

    # Run cluster (bottom-center)
    add_box(s, Inches(4.85), Inches(5.7), Inches(2.0), Inches(0.65),
            "DailyRun\n(id, runDate, status)",
            fill=DANGER, font_size=10)
    add_box(s, Inches(7.0), Inches(5.7), Inches(2.0), Inches(0.65),
            "PaperRunResult\n(finalRank, isRecommended)",
            fill=DANGER, font_size=10)
    # links
    add_arrow(s, Inches(5.85), Inches(5.7),
              Inches(6.0),  Inches(4.3), color=RULE, weight=1.0)
    add_arrow(s, Inches(7.5), Inches(5.7),
              Inches(7.0),  Inches(4.3), color=RULE, weight=1.0)
    add_arrow(s, Inches(6.85), Inches(6.0),
              Inches(7.0),  Inches(6.0), color=RULE, weight=1.0)

    # Note panel
    add_box(s, Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.4),
            "Bilingual JSON columns: summary, keyContribution, "
            "methodologySummary, strengths, weaknesses, rankingExplanation, "
            "recommendationReason, PaperFigure.caption — stored as Prisma `Json?`.",
            fill=BG, text_color=MUTED, font_size=10, bold=False,
            shape=MSO_SHAPE.RECTANGLE)


def slide_19_frontend(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Frontend Surface",
                  "Next.js 16 App Router, server-only repos, i18n cookie")

    # Routes table
    headers = ["Route", "Type", "Backing repo(s)", "Notes"]
    rows = [
        ["/",                       "RSC page", "runsRepo, runResultsRepo, trendsRepo",
         "Home feed (10 / page), sidebar tags, hot topics"],
        ["/library",                "RSC page", "papersRepo",
         "Filterable browse — in progress (controls partly disabled)"],
        ["/papers/[id]",            "RSC page", "papersRepo, evaluationsRepo, sourcesRepo",
         "Single-paper detail (all evals, sources, code links)"],
        ["/runs/[id]",              "RSC page", "runsRepo, runResultsRepo",
         "Per-run summary + ranked paper list"],
        ["/api/locale",             "Route handler", "—",
         "POST sets `locale` cookie (∈ en, zh-TW)"],
        ["/api/runs",               "Route handler", "runsRepo",
         "JSON listing of recent DailyRuns"],
        ["/api/papers/[id]/figure", "Route handler", "figuresRepo",
         "Serves PaperFigure.imageBytes as image/png"],
    ]
    add_table(s, headers, rows,
              left=Inches(0.55), top=Inches(1.55),
              width=Inches(12.3), height=Inches(3.6), font_size=10)

    add_bullets(s, [
        ("Server boundary",
         "src/server/** files are `'server-only'` — accidentally importing "
         "into a client component is a build error. UI talks to data via "
         "the repository layer (src/server/repos/*), never raw Prisma."),
        ("Best-evaluation selection",
         "src/server/lib/select-evaluation.ts picks SUCCESS > ABSTRACT > "
         "FAILED|UNAVAILABLE — same logic shared by ingest ranking and UI display."),
        ("Locale render",
         "src/lib/locale.ts pickLocalized() runs at the component level — "
         "the API ships both locales; the client picks one."),
    ], left=Inches(0.55), top=Inches(5.4), width=Inches(12.3),
       height=Inches(1.6), font_size=11)

    add_footer(s)


def slide_20_testing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Testing Strategy",
                  "Vitest (unit + integration) + Playwright (E2E) + Zod CLI gates")

    headers = ["Tier", "Tool", "Location", "Coverage", "Run"]
    rows = [
        ["Unit",         "Vitest",
         "tests/unit/{sources,dedup,prompt-eval}/*",
         "parseArxivAtom · parseOpenReview · parseHuggingFace · "
         "dedupWithinBatch · normalize · fingerprint · matcher findMatch · "
         "prompt bounds + check-evaluations",
         "npm test"],
        ["Integration",  "Vitest + real Postgres",
         "tests/integration/{ingest,collect-persist}.test.ts",
         "Full ingest of sample run · figure rows · idempotency + --force · "
         "Phase 2.5 reference run · alt-source joinKey · fuzzy-collision · "
         "fail-fast on missing eval · 3-run cross-run dedup (exact + fuzzy)",
         "RUN_INTEGRATION=1\nnpm run test:integration"],
        ["E2E",          "Playwright",
         "tests/e2e/* (config: playwright.config.ts)",
         "Browser walk-throughs of /, /papers/[id], /library, /runs/[id]",
         "npm run test:e2e"],
        ["Schema gate",  "Zod CLI",
         "scripts/validate-{candidates,evaluations}.ts",
         "Parses a run dir against CandidatesFileSchema / "
         "EvaluationsFileSchema and exits 1 on mismatch",
         "npm run validate:candidates ./data/runs/…\nnpm run validate:evaluations ./data/runs/…"],
        ["Prompt gate",  "Custom harness",
         "scripts/prompt-eval/check-evaluations.ts",
         "Schema + per-fixture bounds + ranking coarse-flags (see Slide 15)",
         "npm run prompt:check  ./data/runs/…"],
    ]
    add_table(s, headers, rows,
              left=Inches(0.35), top=Inches(1.55),
              width=Inches(12.6), height=Inches(5.0), font_size=9)

    add_box(s, Inches(0.55), Inches(6.7), Inches(12.3), Inches(0.4),
            "Integration tests use a dedicated test DB (`DATABASE_URL_TEST`); "
            "tests/integration/setup.ts handles setup / cleanup per file.",
            fill=BG, text_color=MUTED, font_size=10, bold=False,
            shape=MSO_SHAPE.RECTANGLE)

    add_footer(s)


def slide_21_risks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, "Open Questions & Risks",
                  "Discussion points for the review")

    headers = ["Area", "Risk / Open question", "Today's mitigation", "Suggested follow-up"]
    rows = [
        ["LLM drift",
         "Subtle bilingual quality regression "
         "(e.g. zh-TW phrasing drift) not detected by the harness.",
         "Phase 2.5 catches gross score drift via per-fixture bounds.",
         "Add per-locale length/character-set assertions; sample human review monthly."],
        ["PDF truncation",
         "Appendix-marker heuristic fails on non-standard headings "
         "(\"Limitations & Discussion\", localized headings).",
         "Falls back to keeping the entire PDF (safe but slower / noisier).",
         "Add a TOC-based truncation when bookmarks present; "
         "track truncation outcome in tableFigureAnalysis."],
        ["Figure extraction",
         "Bounding-box is LLM-estimated; no automatic rendering-quality "
         "check before storing the crop.",
         "Two-pass render at different DPIs catches gross blank crops by feel.",
         "Add OCR-presence check; reject crops below an entropy threshold."],
        ["Figure storage",
         "PaperFigure.imageBytes (BYTEA) inline — fine at ~30/day; "
         "row-size dominates if cadence increases.",
         "Cascade-delete is clean; no orphan blobs.",
         "Move to S3 / R2 with signed URLs; keep figure metadata in DB."],
        ["Auth",
         "No identity layer — PaperFeedback.userId is nullable; UI is anonymous.",
         "Read-only V1 — no destructive writes from the UI.",
         "Decide between Clerk / NextAuth before enabling feedback UI."],
        ["Source coverage",
         "Only 3 sources hard-coded; quotas fixed at 15 / 5 / 10.",
         "Within-batch dedup handles overlap cleanly.",
         "Make source registry pluggable; pull quotas from config."],
    ]
    add_table(s, headers, rows,
              left=Inches(0.35), top=Inches(1.55),
              width=Inches(12.6), height=Inches(5.4), font_size=10)

    add_footer(s)


# ---------- Build ----------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_overview(prs)
    slide_04_stack(prs)
    slide_05_repo(prs)
    slide_06_pipeline(prs)
    slide_07_collection_sources(prs)
    slide_08_aggregation(prs)
    slide_09_candidate_schema(prs)
    slide_10_abstract_screening(prs)
    slide_11_full_pdf(prs)
    slide_12_figure_extraction(prs)
    slide_13_eval_schema(prs)
    slide_14_bilingual(prs)
    slide_15_prompt_harness(prs)
    slide_16_ingest(prs)
    slide_17_dedup(prs)
    slide_18_db_schema(prs)
    slide_19_frontend(prs)
    slide_20_testing(prs)
    slide_21_risks(prs)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
